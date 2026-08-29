#!/usr/bin/env python3
"""gakkai-slides build_deck.py: アウトライン(Markdown)から学会発表用pptxを生成する。

使い方:
    python3 build_deck.py outline.md -o deck.pptx
    python3 build_deck.py outline.md -o deck.pptx --font "ヒラギノ角ゴシック"

依存: python-pptx（pip install python-pptx）。Pillowはpython-pptxと一緒に入る。

アウトラインの書式は references/ と README を参照。要点:
    # 演題名
    演者: 中島誉也、共同演者A、共同演者B
    所属: 長崎大学病院 麻酔集中治療部
    学会: 第XX回 〇〇学会
    日付: 2026年9月X日
    発表時間: 7分
    COI: なし            # 「あり: 内容」/「skip」も可
    比率: 16:9           # 4:3 も可

    ## 背景
    - 箇条書き1
    - **赤太字で強調**を含む行
      - 下位項目（スペース2つでインデント）
    ![](figures/fig1.png)
    主張: このスライドで言いたいこと1行
    出典: 1) Russell JA, et al. N Engl J Med. 2008.

設計方針（referencesのデザイン規則を実装したもの）:
    白背景、黒本文、赤は強調のみ。塗りつぶし図形・角丸カード・アクセント線は作らない。
    1スライドのテキスト枠は タイトル / 本文(1枠) / 主張 / 出典 / ページ番号 の最大5つ。
    タイトル以外の全スライド右下に n / N。フォントは既定で Noto Sans JP。
"""
from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLACK = RGBColor(0x11, 0x11, 0x11)
RED = RGBColor(0xC0, 0x00, 0x00)
GRAY = RGBColor(0x59, 0x59, 0x59)

MARGIN = Inches(0.5)


# ---------------------------------------------------------------------------
# 低レベルヘルパ
# ---------------------------------------------------------------------------

def set_font(run, font_name, size, color=BLACK, bold=False):
    """和文・英数の両方に同じフォントを効かせる。"""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name  # latin
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font_name)


def add_runs(paragraph, text, font_name, size, color=BLACK, bold=False):
    """**...** を赤太字にしてrunを追加する。赤は強調にだけ使う。"""
    for i, part in enumerate(re.split(r"\*\*(.+?)\*\*", text)):
        if part == "":
            continue
        run = paragraph.add_run()
        run.text = part
        if i % 2 == 1:
            set_font(run, font_name, size, RED, True)
        else:
            set_font(run, font_name, size, color, bold)


def set_bullet(paragraph, level, font_name):
    """ぶら下げインデント付きの箇条書きにする。level 0 は「・」、1 は「－」。"""
    pPr = paragraph._p.get_or_add_pPr()
    marL = 285750 * (level + 1)              # 0.25inch * (level+1)
    pPr.set("marL", str(marL))
    pPr.set("indent", str(-285750))
    for tag in ("a:buFont", "a:buChar", "a:buNone", "a:buAutoNum"):
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)
    bf = pPr.makeelement(qn("a:buFont"), {"typeface": font_name})
    pPr.append(bf)
    bu = pPr.makeelement(qn("a:buChar"), {"char": "・" if level == 0 else "－"})
    pPr.append(bu)


def no_bullet(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return box, tf


# ---------------------------------------------------------------------------
# アウトラインの解析
# ---------------------------------------------------------------------------

META_KEYS = {"演者", "所属", "学会", "日付", "発表時間", "COI", "比率"}


def parse_outline(text):
    meta = {"タイトル": "", "比率": "16:9", "COI": "なし"}
    slides = []
    cur = None
    for raw in text.splitlines():
        line = raw.rstrip()
        if not line.strip():
            continue
        if line.startswith("# ") and not meta["タイトル"]:
            meta["タイトル"] = line[2:].strip()
            continue
        m = re.match(r"^(演者|所属|学会|日付|発表時間|COI|比率)\s*[:：]\s*(.+)$", line)
        if m and cur is None:
            meta[m.group(1)] = m.group(2).strip()
            continue
        if line.startswith("## "):
            cur = {"title": line[3:].strip(), "bullets": [], "images": [],
                   "claim": "", "refs": ""}
            slides.append(cur)
            continue
        if cur is None:
            continue
        m = re.match(r"^!\[[^\]]*\]\(([^)]+)\)\s*$", line.strip())
        if m:
            cur["images"].append(m.group(1))
            continue
        m = re.match(r"^(主張|出典)\s*[:：]\s*(.+)$", line.strip())
        if m:
            cur["claim" if m.group(1) == "主張" else "refs"] = m.group(2).strip()
            continue
        m = re.match(r"^(\s*)-\s+(.+)$", line)
        if m:
            level = 1 if len(m.group(1)) >= 2 else 0
            cur["bullets"].append((level, m.group(2).strip()))
            continue
        # 素の行は本文（箇条書きなし）として扱う
        cur["bullets"].append((-1, line.strip()))
    return meta, slides


# ---------------------------------------------------------------------------
# スライドの組み立て
# ---------------------------------------------------------------------------

def build(meta, slides, out_path, font_name):
    prs = Presentation()
    if meta.get("比率", "16:9").replace("：", ":") == "4:3":
        prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
    else:
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    W, H = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def new_slide():
        return prs.slides.add_slide(blank)

    # 1枚目: タイトル
    s = new_slide()
    _, tf = textbox(s, MARGIN, Inches(2.0), W - MARGIN * 2, Inches(2.0))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_runs(p, meta["タイトル"], font_name, 32, BLACK, True)
    _, tf = textbox(s, MARGIN, Inches(4.3), W - MARGIN * 2, Inches(2.2))
    lines = [meta.get("演者", ""), meta.get("所属", ""),
             " ".join(x for x in (meta.get("学会", ""), meta.get("日付", "")) if x)]
    first = True
    for txt in [x for x in lines if x]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)
        add_runs(p, txt, font_name, 20)

    # 2枚目: COI（医学系学会は口演スライド2枚目に開示。学会指定様式があればそちらを優先）
    coi = meta.get("COI", "なし").strip()
    if coi.lower() != "skip":
        s = new_slide()
        _add_title(s, "利益相反（COI）開示", font_name, W)
        _, tf = textbox(s, MARGIN, Inches(2.4), W - MARGIN * 2, Inches(3.0))
        speaker = meta.get("演者", "").split("、")[0].split(",")[0]
        p = tf.paragraphs[0]
        add_runs(p, f"筆頭演者: {speaker}", font_name, 22)
        p = tf.add_paragraph()
        p.space_before = Pt(18)
        if coi in ("なし", "無し", "none"):
            add_runs(p, "本演題に関連し、開示すべきCOI関係にある企業等はありません", font_name, 22)
        else:
            add_runs(p, coi.replace("あり:", "").replace("あり：", "").strip(), font_name, 22)

    # 本文スライド
    for sl in slides:
        s = new_slide()
        _add_title(s, sl["title"], font_name, W)
        body_top = Inches(1.45)
        body_bottom = H - Inches(1.35 if (sl["claim"] or sl["refs"]) else 0.75)
        has_img = bool(sl["images"])
        body_w = (W - MARGIN * 2)
        if has_img and sl["bullets"]:
            body_w = Emu(int((W - MARGIN * 2) * 0.52))

        if sl["bullets"]:
            _, tf = textbox(s, MARGIN, body_top, body_w, body_bottom - body_top)
            first = True
            for level, txt in sl["bullets"]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.space_after = Pt(10)
                if level < 0:
                    no_bullet(p)
                    add_runs(p, txt, font_name, 20)
                else:
                    set_bullet(p, level, font_name)
                    add_runs(p, txt, font_name, 20 if level == 0 else 18,
                             BLACK if level == 0 else BLACK)

        if has_img:
            img_left = MARGIN if not sl["bullets"] else Emu(int(MARGIN + body_w + Inches(0.3)))
            img_w = (W - img_left - MARGIN)
            _place_images(s, sl["images"], img_left, body_top, img_w,
                          body_bottom - body_top)

        if sl["claim"]:
            _, tf = textbox(s, MARGIN, H - Inches(1.25), W - MARGIN * 2, Inches(0.5))
            p = tf.paragraphs[0]
            add_runs(p, sl["claim"], font_name, 20, BLACK, True)
        if sl["refs"]:
            _, tf = textbox(s, MARGIN, H - Inches(0.68), W - Inches(2.0), Inches(0.35))
            p = tf.paragraphs[0]
            add_runs(p, sl["refs"], font_name, 12, GRAY)

    # ページ番号（タイトル以外、右下に n / N）
    total = len(prs.slides.slides if hasattr(prs.slides, "slides") else prs.slides._sldIdLst)
    total = len(prs.slides._sldIdLst)
    for i, s in enumerate(prs.slides, 1):
        if i == 1:
            continue
        _, tf = textbox(s, W - Inches(1.4), H - Inches(0.5), Inches(1.1), Inches(0.35))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        add_runs(p, f"{i} / {total}", font_name, 14, GRAY)

    prs.save(out_path)
    return total


def _add_title(slide, text, font_name, W):
    _, tf = textbox(slide, MARGIN, Inches(0.45), W - MARGIN * 2, Inches(0.85))
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    add_runs(p, text, font_name, 28, BLACK, True)


def _place_images(slide, paths, left, top, width, height):
    from PIL import Image
    n = len(paths)
    cell_h = Emu(int(height / n))
    y = top
    for path in paths:
        if not Path(path).exists():
            print(f"[警告] 画像が見つからない: {path}（スキップ）", file=sys.stderr)
            y = Emu(int(y + cell_h))
            continue
        with Image.open(path) as im:
            iw, ih = im.size
        ratio = ih / iw
        w = width
        h = Emu(int(w * ratio))
        if h > cell_h:
            h = cell_h
            w = Emu(int(h / ratio))
        x = Emu(int(left + (width - w) / 2))
        slide.shapes.add_picture(path, x, y, width=w, height=h)
        y = Emu(int(y + cell_h))


def main():
    ap = argparse.ArgumentParser(description="学会発表スライドをアウトラインから生成する")
    ap.add_argument("outline")
    ap.add_argument("-o", "--out", default="deck.pptx")
    ap.add_argument("--font", default="Noto Sans JP",
                    help="既定は Noto Sans JP。未インストールなら『ヒラギノ角ゴシック』『游ゴシック』等を指定")
    args = ap.parse_args()

    meta, slides = parse_outline(Path(args.outline).read_text(encoding="utf-8"))
    if not meta["タイトル"]:
        sys.exit("エラー: 1行目に「# 演題名」がありません")
    total = build(meta, slides, args.out, args.font)

    print(f"{args.out} を書き出しました（全{total}枚、比率 {meta.get('比率','16:9')}、フォント {args.font}）")
    m = re.search(r"(\d+)", meta.get("発表時間", ""))
    if m:
        minutes = int(m.group(1))
        content = total - (1 if meta.get("COI", "").lower() == "skip" else 2)
        print(f"発表時間 {minutes}分 に対して本文 {content}枚。目安は1分1枚（±2〜3枚）")
    print("次は scripts/check_deck.py で検査してください")


if __name__ == "__main__":
    main()
