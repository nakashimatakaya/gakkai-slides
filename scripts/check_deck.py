#!/usr/bin/env python3
"""gakkai-slides check_deck.py: 学会発表pptxをデザイン規則に照らして検査する。

使い方:
    python3 check_deck.py deck.pptx
    python3 check_deck.py deck.pptx --json

見るもの:
    塗りつぶし図形（box）、角丸カード、ページ番号の欠落、フォント（Noto Sans JP以外）、
    パレット外の色、赤の使いすぎ、16pt未満の本文、テキスト枠の乱立（1文1ボックスの疑い）、
    1枚あたりの行数過多。検出は疑いの提示であって、直すかどうかは人が決める。
依存: python-pptx
"""
from __future__ import annotations

import argparse
import json
import re
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from pptx.oxml.ns import qn

findings = []


def add(level, slide_no, check, message):
    findings.append({"level": level, "slide": slide_no, "check": check, "message": message})


def rgb_of(font):
    try:
        if font.color and font.color.type is not None and font.color.rgb is not None:
            return str(font.color.rgb)
    except Exception:
        pass
    return None


def is_allowed_color(hexstr):
    r, g, b = int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16)
    if abs(r - g) <= 16 and abs(g - b) <= 16:      # 黒〜グレー〜白
        return True
    if r >= 0xA0 and g <= 0x60 and b <= 0x60:      # 赤系
        return True
    return False


def is_red(hexstr):
    r, g, b = int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16)
    return r >= 0xA0 and g <= 0x60 and b <= 0x60


def ea_font(run):
    rPr = run._r.find(qn("a:rPr"))
    if rPr is None:
        return None
    ea = rPr.find(qn("a:ea"))
    return ea.get("typeface") if ea is not None else None


def check(path, font_name):
    prs = Presentation(path)
    n_slides = len(prs.slides._sldIdLst)
    pagenum_re = re.compile(r"^\s*\d+\s*(/\s*\d+)?\s*$")

    W, H = prs.slide_width, prs.slide_height
    placeholder_re = re.compile(r"\bText \d|ラベル\s?\d|タイトル\s?\d|YYYY|Source \d|◯◯|〇〇|（要記入）|\bXX\b")
    titles = []
    all_text = []
    for idx, slide in enumerate(prs.slides, 1):
        # タイトル: 最も上にあるテキスト枠
        tshapes = [sh for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
        if tshapes and idx >= 2:
            top_shape = min(tshapes, key=lambda sh: (sh.top or 0))
            t = top_shape.text_frame.text.strip().splitlines()[0]
            titles.append((idx, t, [sh for sh in tshapes]))
            if re.search(r"(です|ます)[。！]?$", t):
                add("警告", idx, "タイトルの語尾", f"「{t[-6:]}」。体言止めにする")
            if len(t) > 40:
                add("情報", idx, "タイトルが長い", f"{len(t)}字（目安40字以内）")
        for sh in slide.shapes:
            # はみ出し・ページ番号との重なり
            try:
                if sh.left is not None and sh.width is not None and (sh.left + sh.width > W + 12700 or sh.top + sh.height > H + 12700 or sh.left < -12700 or sh.top < -12700):
                    add("警告", idx, "版面からのはみ出し", f"{getattr(sh, 'name', 'shape')} が版面外に出ている")
            except Exception:
                pass
            if sh.has_text_frame:
                txt = sh.text_frame.text
                all_text.append(txt)
                m = placeholder_re.search(txt)
                if m:
                    add("重大", idx, "プレースホルダー残存", f"「{m.group(0)}」が残っている")
        frames = 0
        red_runs = 0
        lines = 0
        has_pagenum = False

        for shape in slide.shapes:
            # 図形（box）の検査
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                add("警告", idx, "図形box", f"オートシェイプ「{shape.shape_type}」がある。boxは使わない")
            if shape.shape_type in (MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.AUTO_SHAPE):
                try:
                    if shape.fill.type is not None and shape.fill.type == 1:  # solid
                        hexstr = str(shape.fill.fore_color.rgb)
                        if hexstr.upper() != "FFFFFF":
                            add("警告", idx, "塗りつぶし", f"塗り #{hexstr} の枠がある。白背景に文字だけにする")
                except Exception:
                    pass
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and "ROUND" in shape.auto_shape_type.name:
                    add("警告", idx, "角丸カード", "角丸図形はAI生成スライドの典型。使わない")
            except Exception:
                pass

            if not shape.has_text_frame:
                continue
            text_all = shape.text_frame.text.strip()
            if text_all:
                frames += 1
            if pagenum_re.match(text_all):
                has_pagenum = True

            for para in shape.text_frame.paragraphs:
                if para.text.strip():
                    lines += 1
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    # フォント
                    for name in (run.font.name, ea_font(run)):
                        if name and not name.startswith(font_name):
                            add("情報", idx, "フォント", f"「{name}」（既定は {font_name}）: {run.text[:14]}")
                            break
                    # サイズ
                    if run.font.size is not None:
                        pt = run.font.size.pt
                        if pt < 12:
                            add("警告", idx, "文字サイズ", f"{pt:.0f}pt（12pt未満）: {run.text[:14]}")
                        elif pt < 16 and not pagenum_re.match(text_all):
                            hexstr0 = rgb_of(run.font)
                            grayish = hexstr0 and abs(int(hexstr0[0:2],16)-int(hexstr0[2:4],16))<=16 and 0x30<=int(hexstr0[0:2],16)<=0x90
                            if not grayish:
                                add("情報", idx, "文字サイズ", f"{pt:.0f}pt。本文は16pt以上、12〜14ptはグレーの出典のみ: {run.text[:14]}")
                    # 色
                    hexstr = rgb_of(run.font)
                    if hexstr:
                        if not is_allowed_color(hexstr):
                            add("警告", idx, "パレット外の色", f"#{hexstr}: {run.text[:14]}。黒・赤・グレーに収める")
                        if is_red(hexstr):
                            red_runs += 1

        if idx >= 2 and not has_pagenum:
            add("警告", idx, "ページ番号", "右下に n / N がない")
        if red_runs > 2:
            add("情報", idx, "赤の使いすぎ", f"赤のrunが{red_runs}箇所。強調は1枚2箇所まで")
        if frames > 6:
            add("情報", idx, "テキスト枠の乱立", f"{frames}枠。1文ごとにboxを分けず、本文は1枠にまとめる")
        if lines > 12:
            add("情報", idx, "行数過多", f"{lines}行。話すことだけ書き、残りは補足スライドへ")

    # タイトル文型の同一率（末尾2文字で近似）
    if len(titles) >= 4:
        ends = [t[-2:] for _, t, _ in titles]
        most = max(set(ends), key=ends.count)
        if ends.count(most) / len(ends) >= 0.6:
            add("情報", 0, "タイトル文型の同一", f"{len(ends)}枚中{ends.count(most)}枚が「…{most}」で終わる。主張でなくテンプレを写していないか")
    # 表記ゆれ
    joined = "\n".join(all_text)
    if re.search(r"[０-９]", joined) and re.search(r"[0-9]", joined):
        add("情報", 0, "表記ゆれ", "全角数字と半角数字が混在")
    if re.search(r"[，．]", joined) and re.search(r"[、。]", joined):
        add("情報", 0, "表記ゆれ", "句読点スタイル（、。と，．）が混在")
    kata = set(re.findall(r"[ァ-ヴ]{3,}ー", joined))
    for w in kata:
        if re.search(re.escape(w[:-1]) + r"(?![ァ-ヴー])", joined):
            add("情報", 0, "表記ゆれ", f"「{w}」と「{w[:-1]}」が混在")
    return n_slides, titles


def print_story(titles, prs_path):
    """--story: タイトルと主張（各スライドで最も下の太字行）の一覧。人が通し読みする。"""
    prs = Presentation(prs_path)
    print(f"== ストーリーライン: {prs_path} ==")
    for idx, slide in enumerate(prs.slides, 1):
        frames = [sh for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
        if not frames:
            continue
        top = min(frames, key=lambda sh: sh.top or 0).text_frame.text.strip().splitlines()[0]
        claim = ""
        for sh in sorted(frames, key=lambda sh: sh.top or 0):
            for p in sh.text_frame.paragraphs:
                if p.runs and all(r.font.bold for r in p.runs if r.text.strip()) and p.text.strip() and sh.top and sh.top > H_OF(prs) * 0.7:
                    claim = p.text.strip()
        print(f"p{idx}  {top}" + (f"  →  {claim}" if claim else ""))


def H_OF(prs):
    return prs.slide_height


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--font", default="Noto Sans")
    ap.add_argument("--json", action="store_true")
    ap.add_argument("--story", action="store_true", help="タイトルと主張の一覧を出して通し読みする")
    args = ap.parse_args()

    n, _titles = check(args.pptx, args.font)
    if args.story:
        print_story(_titles, args.pptx)
        print()
    order = {"重大": 0, "警告": 1, "情報": 2}
    findings.sort(key=lambda f: (order[f["level"]], f["slide"]))

    if args.json:
        print(json.dumps({"slides": n, "findings": findings}, ensure_ascii=False, indent=2))
        return
    counts = {k: sum(1 for f in findings if f["level"] == k) for k in order}
    print(f"== {args.pptx} == 全{n}枚 / 重大 {counts['重大']} / 警告 {counts['警告']} / 情報 {counts['情報']}")
    for f in findings:
        print(f"[{f['level']}] p{f['slide']} {f['check']}: {f['message']}")
    if not findings:
        print("検出なし。")


if __name__ == "__main__":
    main()
